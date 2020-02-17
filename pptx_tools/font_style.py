"""
This module provides a helper class to deal with fonts in python-pptx.
@author: Nathanael Jöhrmann
"""
from typing import Union, Optional, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.enum.text import MSO_TEXT_UNDERLINE_TYPE
from pptx.shapes.autoshape import Shape
from pptx.text.text import Font
from pptx.text.text import _Paragraph
from pptx.text.text import _Run
from pptx.util import Pt

from pptx_tools.fill_style import PPTXFillStyle
from pptx_tools.utils import _USE_DEFAULT


class _DO_NOT_CHANGE:
    def __str__(self):
        return """used to tell PPTXFontStyle.set() to not change a value"""


class PPTXFontStyle:
    """
    Helper class to deal with fonts in python-pptx. The internal class pptx.text.text.Font is limited, as it
    always needs an existing Text/Character/... for initializing and also basic functionality like assignment
    of one font to another is missing.
    """
    # default language and font; no _USE_DEFAULT for language_id -> use MSO_LANGUAGE_ID.NONE
    language_id: Union[MSO_LANGUAGE_ID, _USE_DEFAULT, None] = MSO_LANGUAGE_ID.ENGLISH_UK  # MSO_LANGUAGE_ID.GERMAN
    name: Union[str, _USE_DEFAULT, None] = "Roboto"  # "Arial"  # "Arial Narrow"

    def __init__(self):
        #  If set to use_default(), the bold, italic ... setting is cleared and is inherited
        #  from an enclosing shape’s setting, or a setting in a style or master
        self.bold: Union[bool, _USE_DEFAULT, None] = None
        self.italic: Optional[bool, _USE_DEFAULT, None] = None
        self.underline: Union[MSO_TEXT_UNDERLINE_TYPE, _USE_DEFAULT, bool, None] = None

        # use class attribute; instance attribute only when changed by user
        # self.language_id: MSO_LANGUAGE_ID = MSO_LANGUAGE_ID.NONE  # ENGLISH_UK; ENGLISH_US; ESTONIAN; GERMAN; ...
        # self.name: Union[str, _USE_DEFAULT, None] = None

        # saved in units of Pt (not EMU like pptx.text.text.Font) - converting to EMU is done during write_to_font
        self.size: Optional[int] = None  # 18


        # todo: color is ColorFormat object
        self._color_rgb: Optional[RGBColor] = None
        # fil.fore_color changes font color; also gradient or image might be useful (not implemented in FillStyle jet)
        self.fill_style: Optional[PPTXFillStyle] = None  # PPTXFillStyle()

    @property
    def color_rgb(self):
        return self._color_rgb

    @color_rgb.setter
    def color_rgb(self, value: Union[RGBColor, Tuple[any, any, any], None]):
        assert isinstance(value, RGBColor) or isinstance(value, tuple) or (value is None)
        if isinstance(value, tuple):
            self._color_rgb = RGBColor(*value)
        else:
            self._color_rgb = value


    def read_font(self, font: Font) -> 'PPTXFontStyle':  # todo: check for None behavior (use_dfault() ? )
        """Read attributes from a pptx.text.text.Font object."""
        self.bold = font.bold
        self.italic = font.italic
        self.name = font.name
        self.size = font.size.pt  # todo: convert to Pt ?
        self.underline = font.underline
        return self

    def write_font(self, font: Font) -> None:
        """Write attributes to a pptx.text.text.Font object."""
        font.name = self._get_write_value(new_value=self.name, old_value=font.name)
        font.bold = self._get_write_value(new_value=self.bold, old_value=font.bold)
        font.italic = self._get_write_value(new_value=self.italic, old_value=font.italic)
        font.underline = self._get_write_value(new_value=self.underline, old_value=font.underline)

        if self.language_id == _USE_DEFAULT:
            font.language_id = MSO_LANGUAGE_ID.NONE
        else:
            font.language_id = self._get_write_value(new_value=self.language_id, old_value=font.language_id)

        if self.size is not None:
            if self.size == _USE_DEFAULT:
                font.size = None
            else:
                font.size = Pt(self.size)

        if self.color_rgb is not None:
            font.color.rgb = self.color_rgb

        if self.fill_style is not None:
            self.fill_style.write_fill(font.fill)

    def _get_write_value(self, new_value, old_value, check_default=True):
        """Used to check for None and use_default(), returning the correct value to write."""
        if new_value is None:
            return old_value
        if check_default and (new_value == _USE_DEFAULT):
            return None
        return new_value

    def write_shape(self, shape: Shape) -> None:  # todo: remove? better use write_text_fame
        """
        Write attributes to all paragraphs in given pptx.shapes.autoshape.Shape.
        Raises TypeError if given shape has no text_frame.
        """
        if not shape.has_text_frame:
            raise TypeError("Cannot write font for given shape (has no text_frame)")
        self.write_text_frame(shape.text_frame)

    def write_text_frame(self, text_frame):
        for paragraph in text_frame.paragraphs:
            self.write_paragraph(paragraph)

    def write_paragraph(self, paragraph: _Paragraph) -> None:
        """ Write attributes to given paragraph"""
        self.write_font(paragraph.font)

    def write_run(self, run: _Run) -> None:
        """ Write attributes to given run"""
        self.write_font(run.font)

    @classmethod
    def copy_font(cls, _from: Font, _to: Font) -> None:
        """Copies settings from one pptx.text.text.Font to another."""
        font_style=cls()
        font_style.read_font(_from)
        font_style.write_font(_to)
        # _to.bold = _from.bold
        # # todo: color is ColorFormat object
        # # _to.set_color = _from.color
        # # todo: fill is FillFormat object
        # # _to.fill = _from.fill
        # _to.italic = _from.italic
        # _to.language_id = _from.language_id
        # _to.name = _from.name
        # _to.size = _from.size
        # _to.underline = _from.underline

    def set(self, bold: Optional[bool] = _DO_NOT_CHANGE,
            italic: Optional[bool] = _DO_NOT_CHANGE,
            language_id: MSO_LANGUAGE_ID = _DO_NOT_CHANGE,
            name: Optional[str] = _DO_NOT_CHANGE,
            size: Optional[int] = _DO_NOT_CHANGE,
            underline: Union[MSO_TEXT_UNDERLINE_TYPE, bool, None] = _DO_NOT_CHANGE,
            color_rgb: Union[RGBColor, Tuple[any, any, any]] = _DO_NOT_CHANGE
            ) -> 'PPTXFontStyle':
        """Convenience method to set several font attributes together."""
        if bold is not _DO_NOT_CHANGE:
            self.bold = bold
        if italic is not _DO_NOT_CHANGE:
            self.italic = italic
        if language_id is not _DO_NOT_CHANGE:
            self.language_id = language_id
        if name is not _DO_NOT_CHANGE:
            self.name = name
        if size is not _DO_NOT_CHANGE:
            self.size = size
        if underline is not _DO_NOT_CHANGE:
            self.underline = underline
        if color_rgb is not _DO_NOT_CHANGE:
            self.color_rgb = color_rgb
        return self

        # -----------------------------------------------------------------------------------------------
        # ----------------------------------- experimentell methods -------------------------------------
        # -----------------------------------------------------------------------------------------------
    def _write_font_experimentell(self, font: Font,all_caps: bool = True, strikethrough: bool = True):
        if all_caps:
            font._element.attrib['cap'] = "all"
        else:
            pass

        if strikethrough:
            font._element.attrib['strike'] = "sngStrike"
        else:
            pass
